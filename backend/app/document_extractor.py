"""
Extractores de documentos.

Todos los comentarios y docstrings usan espanol sin tildes.
"""

from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation


def extract_pdf(path: Path) -> list[dict]:
    """
    Extrae paginas PDF.
    """

    reader = PdfReader(str(path))

    output = []

    for index, page in enumerate(reader.pages):

        text = page.extract_text() or ""

        output.append(
            {
                "page": index + 1,
                "text": text,
            }
        )

    return output


def extract_docx(path: Path) -> list[dict]:
    """
    Extrae texto DOCX.
    """

    document = Document(str(path))

    paragraphs = []

    for paragraph in document.paragraphs:
        paragraphs.append(paragraph.text)

    return [
        {
            "page": "",
            "text": "\n".join(paragraphs),
        }
    ]


def extract_xlsx(path: Path) -> list[dict]:
    """
    Extrae contenido XLSX.
    """

    workbook = load_workbook(
        filename=str(path),
        data_only=True,
    )

    sheets = []

    for sheet in workbook.worksheets:

        rows = []

        for row in sheet.iter_rows(values_only=True):

            values = []

            for value in row:

                if value is None:
                    continue

                values.append(str(value))

            if values:
                rows.append(" | ".join(values))

        sheets.append(
            {
                "page": sheet.title,
                "text": "\n".join(rows),
            }
        )

    return sheets


def extract_pptx(path: Path) -> list[dict]:
    """
    Extrae contenido PPTX.
    """

    presentation = Presentation(str(path))

    slides = []

    for slide_index, slide in enumerate(presentation.slides):

        texts = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                texts.append(shape.text)

        slides.append(
            {
                "page": slide_index + 1,
                "text": "\n".join(texts),
            }
        )

    return slides